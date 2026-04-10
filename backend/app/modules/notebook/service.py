"""
Notebook service - handles AI chat, document processing, voice, and studio actions.
"""

import os
import subprocess
import tempfile
import uuid
import httpx
from typing import Optional, List, Dict, Any
from datetime import datetime
from loguru import logger

from app.core.config import settings
from app.odm.notebook import (
    NotebookSourceDocument,
    ChatMessageModel,
    NotebookChatDocument,
    NotebookSettingsDocument,
    NotebookNoteDocument,
)


# ===================== Upload Service =====================

async def upload_source(user_id: str, file_content: bytes, filename: str, content_type: str) -> NotebookSourceDocument:
    """Upload and store a document source."""
    # Determine file type
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else "unknown"
    file_type_map = {
        "pdf": "pdf", "doc": "docx", "docx": "docx",
        "ppt": "pptx", "pptx": "pptx",
        "jpg": "image", "jpeg": "image", "png": "image", "gif": "image",
        "webp": "image", "bmp": "image",
    }
    file_type = file_type_map.get(ext, "other")

    # Save file
    unique_name = f"{uuid.uuid4().hex}_{filename}"
    upload_dir = "app/static/notebooks"
    os.makedirs(upload_dir, exist_ok=True)
    file_path = os.path.join(upload_dir, unique_name)

    with open(file_path, "wb") as f:
        f.write(file_content)

    file_url = f"/static/notebooks/{unique_name}"

    source = await NotebookSourceDocument.create(
        user_id=user_id,
        filename=unique_name,
        original_filename=filename,
        file_url=file_url,
        file_type=file_type,
        file_size=len(file_content),
        processing_status="uploaded",
    )
    return source


async def get_sources(user_id: str) -> List[NotebookSourceDocument]:
    """Get all sources for a user."""
    return await NotebookSourceDocument.find(
        {"user_id": user_id},
        sort=[("created_at", -1)],
        limit=100,
    )


async def delete_source(user_id: str, source_id: str) -> bool:
    """Delete a source."""
    source = await NotebookSourceDocument.find_by_id(source_id)
    if not source or source.user_id != user_id:
        return False

    # Delete file from disk
    file_path = f"app{source.file_url}" if source.file_url.startswith("/static") else source.file_url
    if os.path.exists(file_path):
        os.remove(file_path)

    return await source.delete()


async def extract_source_text(source_id: str) -> Optional[str]:
    """Extract text from a source document (basic extraction)."""
    source = await NotebookSourceDocument.find_by_id(source_id)
    if not source:
        return None

    if source.extracted_text:
        return source.extracted_text

    file_path = f"app{source.file_url}" if source.file_url.startswith("/static") else source.file_url

    extracted = ""
    try:
        if source.file_type == "pdf":
            # Basic PDF text extraction
            try:
                import fitz  # PyMuPDF
                doc = fitz.open(file_path)
                for page in doc:
                    extracted += page.get_text() + "\n"
                doc.close()
            except ImportError:
                extracted = "[PDF text extraction requires PyMuPDF. Install with: pip install PyMuPDF]"

        elif source.file_type == "docx":
            try:
                import docx
                doc = docx.Document(file_path)
                for para in doc.paragraphs:
                    extracted += para.text + "\n"
            except ImportError:
                extracted = "[DOCX extraction requires python-docx. Install with: pip install python-docx]"

        elif source.file_type == "pptx":
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            extracted += shape.text + "\n"
            except ImportError:
                extracted = "[PPTX extraction requires python-pptx. Install with: pip install python-pptx]"

        elif source.file_type == "image":
            # For images, we'll use the AI vision model later
            extracted = "[Image - will use AI vision for extraction]"

        # Save extracted text
        if extracted:
            source.extracted_text = extracted
            source.processing_status = "ready"
            await source.save()

    except Exception as e:
        logger.error(f"Error extracting text from source {source_id}: {e}")
        source.processing_status = "error"
        await source.save()

    return extracted


# ===================== AI Service =====================

async def get_global_settings() -> NotebookSettingsDocument:
    """Get or create global AI settings (admin-configured)."""
    s = await NotebookSettingsDocument.find_one({"settings_key": "global"})
    if not s:
        s = await NotebookSettingsDocument.create(settings_key="global")
    return s


async def update_global_settings(updates: Dict[str, Any]) -> NotebookSettingsDocument:
    """Update global AI settings (admin only)."""
    s = await get_global_settings()
    for key, value in updates.items():
        if value is not None and hasattr(s, key):
            setattr(s, key, value)
    await s.save()
    return s


async def test_connection() -> Dict[str, Any]:
    """Test the AI API connection."""
    s = await get_global_settings()

    if not s.api_base_url or not s.api_key:
        return {"success": False, "message": "API credentials not configured."}

    try:
        async with httpx.AsyncClient(timeout=15.0) as client:
            api_url = s.api_base_url.rstrip("/")
            if not api_url.endswith("/chat/completions"):
                api_url += "/chat/completions"

            response = await client.post(
                api_url,
                headers={
                    "Authorization": f"Bearer {s.api_key}",
                    "Content-Type": "application/json",
                },
                json={
                    "model": s.model_knowledge_chat,
                    "messages": [{"role": "user", "content": "ping"}],
                    "max_tokens": 5,
                },
            )

            if response.status_code == 200:
                return {"success": True, "message": "Connection successful!", "model": s.model_knowledge_chat}
            else:
                return {"success": False, "message": f"API returned status {response.status_code}: {response.text[:150]}"}

    except httpx.ConnectError:
        return {"success": False, "message": "Connection error. Check the API Base URL."}
    except httpx.TimeoutException:
        return {"success": False, "message": "Connection timed out. The API server may be down."}
    except Exception as e:
        return {"success": False, "message": f"Connection failed: {str(e)}"}


async def send_ai_message(
    user_id: str,
    message: str,
    chat_id: Optional[str] = None,
    source_ids: Optional[List[str]] = None,
) -> Dict[str, Any]:
    """Send a message to the AI and get a response."""
    global_settings = await get_global_settings()

    if not global_settings.api_base_url or not global_settings.api_key:
        return {
            "error": True,
            "content": "AI API Key not configured. Please add it to your Admin Settings.",
            "chat_id": chat_id,
        }

    # Get or create chat
    chat = None
    if chat_id:
        chat = await NotebookChatDocument.find_by_id(chat_id)

    if not chat:
        # Create new chat with first message as title
        title = message[:50] + "..." if len(message) > 50 else message
        chat = await NotebookChatDocument.create(
            user_id=user_id,
            title=title,
            source_ids=source_ids or [],
        )

    chat_id_str = str(chat.id)

    # Save user message
    user_msg = await ChatMessageModel.create(
        chat_id=chat_id_str,
        user_id=user_id,
        role="user",
        content=message,
    )

    # Build context from sources
    context_text = ""
    effective_source_ids = source_ids or chat.source_ids
    if effective_source_ids:
        for sid in effective_source_ids:
            text = await extract_source_text(sid)
            if text and not text.startswith("["):
                context_text += f"\n--- Source Document ---\n{text[:3000]}\n"

    # Build messages for AI
    history = await ChatMessageModel.find(
        {"chat_id": chat_id_str},
        sort=[("created_at", 1)],
        limit=20,
    )

    messages = []
    system_prompt = "You are a helpful AI assistant that helps users understand and analyze their documents. Be concise, accurate, and helpful."
    if context_text:
        system_prompt += f"\n\nHere are the relevant documents for context:\n{context_text}"

    messages.append({"role": "system", "content": system_prompt})

    for msg in history:
        messages.append({"role": msg.role, "content": msg.content})

    # Call AI API
    try:
        async with httpx.AsyncClient(timeout=60.0) as client:
            api_url = global_settings.api_base_url.rstrip("/")
            if not api_url.endswith("/chat/completions"):
                api_url += "/chat/completions"

            response = await client.post(
                api_url,
                headers={
                    "Authorization": f"Bearer {global_settings.api_key}",
                    "Content-Type": "application/json",
                },
                json={
                    "model": global_settings.model_knowledge_chat,
                    "messages": messages,
                    "temperature": global_settings.temperature,
                    "max_tokens": global_settings.max_tokens,
                },
            )

            if response.status_code != 200:
                error_detail = response.text[:200]
                logger.error(f"AI API error: {response.status_code} - {error_detail}")

                # Save error as assistant message
                error_msg = await ChatMessageModel.create(
                    chat_id=chat_id_str,
                    user_id=user_id,
                    role="assistant",
                    content=f"Error: Request failed with status code {response.status_code}\n\nThis usually happens when the AI service is overloaded or the API quota is reached. Please try again in a few moments.",
                    metadata={"error": True},
                )

                return {
                    "error": True,
                    "content": error_msg.content,
                    "message_id": str(error_msg.id),
                    "chat_id": chat_id_str,
                    "chat_title": chat.title,
                }

            data = response.json()
            ai_content = data["choices"][0]["message"]["content"]

    except Exception as e:
        logger.error(f"AI API exception: {e}")
        ai_content = f"Error: Could not connect to AI service. Please check your API settings.\n\nDetails: {str(e)}"

    # Save assistant message
    assistant_msg = await ChatMessageModel.create(
        chat_id=chat_id_str,
        user_id=user_id,
        role="assistant",
        content=ai_content,
    )

    return {
        "error": False,
        "content": ai_content,
        "message_id": str(assistant_msg.id),
        "chat_id": chat_id_str,
        "chat_title": chat.title,
    }


# ===================== Chat Management =====================

async def get_chats(user_id: str) -> List[NotebookChatDocument]:
    """Get all chats for a user."""
    return await NotebookChatDocument.find(
        {"user_id": user_id, "is_active": True},
        sort=[("updated_at", -1)],
        limit=50,
    )


async def get_chat_with_messages(user_id: str, chat_id: str) -> Optional[Dict]:
    """Get a chat with its messages."""
    chat = await NotebookChatDocument.find_by_id(chat_id)
    if not chat or chat.user_id != user_id:
        return None

    messages = await ChatMessageModel.find(
        {"chat_id": str(chat.id)},
        sort=[("created_at", 1)],
        limit=200,
    )

    return {"chat": chat, "messages": messages}


async def rename_chat(user_id: str, chat_id: str, title: str) -> bool:
    """Rename a chat."""
    chat = await NotebookChatDocument.find_by_id(chat_id)
    if not chat or chat.user_id != user_id:
        return False
    chat.title = title
    await chat.save()
    return True


async def delete_chat(user_id: str, chat_id: str) -> bool:
    """Delete a chat and its messages."""
    chat = await NotebookChatDocument.find_by_id(chat_id)
    if not chat or chat.user_id != user_id:
        return False

    # Delete all messages in this chat
    await ChatMessageModel.delete_many({"chat_id": str(chat.id)})
    return await chat.delete()


# ===================== Voice Service =====================

def convert_audio_to_wav(audio_data: bytes, original_ext: str = "webm") -> bytes:
    """Uses FFmpeg to transcode incoming audio into WAV for strict APIs."""
    with tempfile.TemporaryDirectory() as td:
        in_path = os.path.join(td, f"input.{original_ext}")
        out_path = os.path.join(td, "output.wav")
        with open(in_path, "wb") as f:
            f.write(audio_data)
            
        cmd = ["ffmpeg", "-y", "-i", in_path, "-vn", "-acodec", "pcm_s16le", "-ar", "16000", "-ac", "1", out_path]
        try:
            subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
            with open(out_path, "rb") as f:
                return f.read()
        except subprocess.CalledProcessError as e:
            err_msg = e.stderr.decode('utf-8', errors='ignore')
            logger.error(f"FFmpeg conversion failed: {err_msg}")
            raise ValueError("Failed to convert audio format using FFmpeg.")


async def speech_to_text(audio_data: bytes, filename: str) -> str:
    """Convert speech to text using Whisper API."""
    whisper_base = getattr(settings, "WHISPER_API_BASE", "")
    whisper_key = getattr(settings, "WHISPER_API_KEY", "")

    if not whisper_base or not whisper_key:
        raise ValueError("Whisper API not configured. Set WHISPER_API_BASE and WHISPER_API_KEY in environment.")

    if not audio_data or len(audio_data) < 100:
        raise ValueError("Audio data is empty or too small.")

    # Convert audio to WAV if it's not strictly mp3/wav already
    ext = filename.rsplit(".", 1)[-1].lower() if "." in filename else "webm"
    if ext not in ["wav", "mp3"]:
        audio_data = convert_audio_to_wav(audio_data, ext)
        ext = "wav"
        filename = "audio.wav"

    # Determine content type from filename extension
    content_type_map = {
        "wav": "audio/wav",
        "mp3": "audio/mpeg",
    }
    content_type = content_type_map.get(ext, "audio/wav")

    # Build the transcriptions URL correctly.
    api_url = whisper_base.rstrip("/")
    if not api_url.endswith("/audio/transcriptions"):
        if api_url.endswith("/v1"):
            api_url += "/audio/transcriptions"
        elif api_url.endswith("/v1/audio"):
            api_url += "/transcriptions"
        else:
            api_url += "/audio/transcriptions"

    # Try the constructed URL first, then fallback with /v1 prefix
    urls_to_try = [api_url]
    # Also try /v1/audio/transcriptions if our URL doesn't have /v1
    if "/v1/" not in api_url:
        base = whisper_base.rstrip("/")
        urls_to_try.append(f"{base}/v1/audio/transcriptions")

    last_error = None
    for url in urls_to_try:
        try:
            async with httpx.AsyncClient(timeout=60.0) as client:
                response = await client.post(
                    url,
                    headers={"Authorization": f"Bearer {whisper_key}"},
                    files={"file": (filename, audio_data, content_type)},
                    data={"model": "whisper-large-v3"},
                )

                if response.status_code == 200:
                    data = response.json()
                    text = data.get("text", "").strip()
                    return text

                logger.warning(f"Whisper API ({url}) returned {response.status_code}: {response.text[:150]}")
                last_error = f"Whisper API returned status {response.status_code}"

        except httpx.TimeoutException:
            logger.error(f"Whisper API request timed out for {url}")
            last_error = "Speech-to-text request timed out. Please try a shorter recording."
        except httpx.RequestError as e:
            logger.error(f"Whisper API connection error for {url}: {e}")
            last_error = f"Could not connect to Whisper API: {e}"

    raise ValueError(last_error or "Speech-to-text failed.")


async def _tts_elevenlabs(text: str, voice_id: str) -> bytes:
    """Text-to-speech via ElevenLabs API."""
    api_key = getattr(settings, "ELEVENLABS_API_KEY", "")
    if not api_key:
        raise ValueError("ElevenLabs API not configured.")

    # Use a known working premade voice if default is passed
    if voice_id == "21m00Tcm4TlvDq8ikWAM":
        voice_id = "CwhRBWXzGAHq8TQ4Fs17"  # Roger - Laid-Back, Casual

    url = f"https://api.elevenlabs.io/v1/text-to-speech/{voice_id}"

    async with httpx.AsyncClient(timeout=120.0) as client:
        response = await client.post(
            url,
            headers={
                "xi-api-key": api_key,
                "Content-Type": "application/json",
                "Accept": "audio/mpeg",
            },
            json={
                "text": text[:5000],
                "model_id": "eleven_flash_v2_5",
                "voice_settings": {
                    "stability": 0.5,
                    "similarity_boost": 0.75,
                },
            },
        )

        if response.status_code == 401:
            raise ValueError("ElevenLabs API authentication failed. Check your API key.")
        elif response.status_code == 402:
            raise ValueError("ElevenLabs account has insufficient credits. Please add credits at elevenlabs.io.")
        elif response.status_code == 429:
            raise ValueError("ElevenLabs rate limit exceeded. Please try again later.")
        elif response.status_code != 200:
            logger.error(f"ElevenLabs API error: {response.status_code} - {response.text[:200]}")
            raise ValueError(f"ElevenLabs API returned status {response.status_code}")

        content_type = response.headers.get("content-type", "")
        if "audio" not in content_type and "octet-stream" not in content_type:
            raise ValueError("ElevenLabs API returned non-audio response.")

        return response.content


async def _tts_openai_compatible(text: str, voice: str = "alloy") -> bytes:
    """Text-to-speech via OpenAI-compatible /v1/audio/speech endpoint.

    Tries these API sources in order:
    1. PODCAST_API_BASE / PODCAST_API_KEY (dedicated podcast TTS)
    2. WHISPER_API_BASE / WHISPER_API_KEY (often same OpenAI-compatible server)
    3. Notebook global settings AI API
    """
    candidates = [
        (getattr(settings, "PODCAST_API_BASE", ""), getattr(settings, "PODCAST_API_KEY", "")),
        (getattr(settings, "WHISPER_API_BASE", ""), getattr(settings, "WHISPER_API_KEY", "")),
    ]

    # Also try notebook global settings
    global_settings = await get_global_settings()
    if global_settings.api_base_url and global_settings.api_key:
        candidates.append((global_settings.api_base_url, global_settings.api_key))

    # Filter out empty entries
    candidates = [(b, k) for b, k in candidates if b and k]

    if not candidates:
        raise ValueError("No TTS service configured. Set ELEVENLABS_API_KEY, PODCAST_API_BASE, or configure AI API settings.")

    last_error = None
    for api_base, api_key in candidates:
        # Build the TTS URL
        api_url = api_base.rstrip("/")
        # Strip known suffixes to get the base
        for suffix in ["/chat/completions", "/v1/chat/completions", "/v1/audio/transcriptions", "/v1/audio"]:
            if api_url.endswith(suffix):
                api_url = api_url[: -len(suffix)]
                break
        # Ensure it ends with /v1 base so we can append /audio/speech
        if api_url.endswith("/v1"):
            api_url += "/audio/speech"
        elif not api_url.endswith("/v1/audio/speech"):
            api_url = api_url.rstrip("/") + "/v1/audio/speech"

        try:
            async with httpx.AsyncClient(timeout=120.0) as client:
                response = await client.post(
                    api_url,
                    headers={
                        "Authorization": f"Bearer {api_key}",
                        "Content-Type": "application/json",
                    },
                    json={
                        "model": "tts-1",
                        "input": text[:4096],
                        "voice": voice,
                        "response_format": "mp3",
                    },
                )

                if response.status_code == 200:
                    if len(response.content) > 0:
                        return response.content
                    raise ValueError("Empty audio response")

                logger.warning(f"OpenAI TTS ({api_url}) returned {response.status_code}: {response.text[:150]}")
                last_error = f"TTS API returned status {response.status_code}"

        except httpx.TimeoutException:
            last_error = "TTS request timed out"
            logger.warning(f"TTS timeout for {api_url}")
        except httpx.RequestError as e:
            last_error = f"TTS connection error: {e}"
            logger.warning(f"TTS connection error for {api_url}: {e}")
        except ValueError as e:
            last_error = str(e)

    raise ValueError(f"All TTS endpoints failed. Last error: {last_error}")


async def _tts_google(text: str) -> bytes:
    """Text-to-speech via Google Cloud Text-to-Speech API."""
    api_key = getattr(settings, "GOOGLE_API_KEY", "")
    if not api_key:
        raise ValueError("Google API key not configured.")

    url = f"https://texttospeech.googleapis.com/v1/text:synthesize?key={api_key}"

    # Truncate to ~5000 chars (Google TTS limit is 5000 bytes of input)
    input_text = text[:5000]

    async with httpx.AsyncClient(timeout=120.0) as client:
        response = await client.post(
            url,
            json={
                "input": {"text": input_text},
                "voice": {
                    "languageCode": "en-US",
                    "name": "en-US-Neural2-J",
                    "ssmlGender": "MALE",
                },
                "audioConfig": {
                    "audioEncoding": "MP3",
                    "speakingRate": 1.0,
                    "pitch": 0.0,
                },
            },
        )

        if response.status_code != 200:
            logger.warning(f"Google TTS error: {response.status_code} - {response.text[:200]}")
            raise ValueError(f"Google TTS returned status {response.status_code}")

        data = response.json()
        audio_content = data.get("audioContent", "")
        if not audio_content:
            raise ValueError("Google TTS returned empty audio.")

        import base64
        return base64.b64decode(audio_content)


async def text_to_speech(text: str, voice_id: str = "21m00Tcm4TlvDq8ikWAM") -> bytes:
    """Convert text to speech. Tries multiple providers in order:
    1. ElevenLabs (if ELEVENLABS_API_KEY set)
    2. OpenAI-compatible (PODCAST_API_BASE, WHISPER_API_BASE, or AI settings)
    3. Google Cloud TTS (if GOOGLE_API_KEY set)
    """
    if not text or not text.strip():
        raise ValueError("Text cannot be empty.")

    errors = []

    # 1. Try ElevenLabs
    elevenlabs_key = getattr(settings, "ELEVENLABS_API_KEY", "")
    if elevenlabs_key:
        try:
            return await _tts_elevenlabs(text, voice_id)
        except Exception as e:
            logger.warning(f"ElevenLabs TTS failed: {e}")
            errors.append(f"ElevenLabs: {e}")

    # 2. Try OpenAI-compatible (PODCAST_API_BASE, WHISPER_API_BASE, etc.)
    try:
        return await _tts_openai_compatible(text)
    except Exception as e:
        logger.warning(f"OpenAI-compatible TTS failed: {e}")
        errors.append(f"OpenAI-compatible: {e}")

    # 3. Try Google Cloud TTS
    google_key = getattr(settings, "GOOGLE_API_KEY", "")
    if google_key:
        try:
            return await _tts_google(text)
        except Exception as e:
            logger.warning(f"Google TTS failed: {e}")
            errors.append(f"Google: {e}")

    logger.error(f"All TTS methods failed: {errors}")
    raise ValueError(f"Text-to-speech failed. Please configure a TTS service (ElevenLabs, OpenAI-compatible, or Google Cloud TTS).")


# ===================== Studio Service =====================

STUDIO_PROMPTS = {
    "summary": "Provide a comprehensive summary of the following documents. Organize it with clear headings and bullet points.",
    "quiz": "Generate a quiz with 10 multiple-choice questions based on the following documents. Format each question with options A-D and indicate the correct answer.",
    "faq": "Generate a list of 10 frequently asked questions with detailed answers based on the following documents.",
    "infographic": """Generate structured infographic content based on the following documents. Return a JSON object with this structure:
{
  "title": "Main title",
  "sections": [
    {
      "heading": "Section heading",
      "key_points": ["point 1", "point 2"],
      "icon": "suggested-icon-name"
    }
  ],
  "statistics": [{"label": "stat label", "value": "stat value"}]
}""",
    "flashcards": "Generate 15 flashcards based on the following documents. Format each as:\n**Front:** [question/term]\n**Back:** [answer/definition]",
    "study-guide": "Create a detailed study guide from the following documents. Include key concepts, definitions, important facts, and study tips organized by topic.",
    "briefing": "Create a professional briefing document from the following documents. Include executive summary, key findings, action items, and recommendations.",
    "exam": "Create a comprehensive exam with 20 questions covering the following documents. Include a mix of multiple choice, true/false, short answer, and essay questions. Provide an answer key at the end.",
    "podcast": """Generate a natural, engaging podcast conversation between two hosts discussing the following documents.
The hosts are:
- Alex: The main host who guides the conversation, asks insightful questions, and keeps things on track.
- Sam: The expert co-host who dives deep into the details and provides analysis.

Guidelines:
- Make it sound like a real podcast conversation with natural dialogue, reactions, and transitions.
- Start with a brief intro where Alex welcomes listeners and introduces the topic.
- Cover all the key points from the documents in an engaging, conversational way.
- Include moments of surprise, humor, and genuine curiosity.
- Use phrases like "That's a great point", "What really caught my eye was...", "Let me break that down..."
- End with a summary of key takeaways and a sign-off.
- The conversation should be 800-1500 words long for a good audio experience.
- Do NOT use any markdown formatting, headers, or bullet points - write it as pure dialogue.
- Format each line as: "Alex: ..." or "Sam: ..."

Write ONLY the dialogue, nothing else.""",
}


async def run_studio_action(user_id: str, action: str, source_ids: List[str]) -> Dict[str, Any]:
    """Run a studio action (summary, quiz, etc.) on sources."""
    global_settings = await get_global_settings()

    if not global_settings.api_base_url or not global_settings.api_key:
        return {
            "error": True,
            "content": "AI API Key not configured. Please add it to your Admin Settings.",
        }

    prompt = STUDIO_PROMPTS.get(action)
    if not prompt:
        return {"error": True, "content": f"Unknown action: {action}"}

    # Gather source texts
    context_text = ""
    for sid in source_ids:
        text = await extract_source_text(sid)
        if text and not text.startswith("["):
            context_text += f"\n--- Document ---\n{text[:4000]}\n"

    if not context_text:
        return {"error": True, "content": "No document content available. Please upload sources first."}

    full_prompt = f"{prompt}\n\n{context_text}"

    try:
        async with httpx.AsyncClient(timeout=90.0) as client:
            api_url = global_settings.api_base_url.rstrip("/")
            if not api_url.endswith("/chat/completions"):
                api_url += "/chat/completions"

            response = await client.post(
                api_url,
                headers={
                    "Authorization": f"Bearer {global_settings.api_key}",
                    "Content-Type": "application/json",
                },
                json={
                    "model": global_settings.model_knowledge_chat,
                    "messages": [
                        {"role": "system", "content": "You are an expert educational content creator. Generate high-quality, accurate content based on the provided documents."},
                        {"role": "user", "content": full_prompt},
                    ],
                    "temperature": 0.7,
                    "max_tokens": 4096,
                },
            )

            if response.status_code != 200:
                return {"error": True, "content": f"AI API error: status {response.status_code}"}

            data = response.json()
            content = data["choices"][0]["message"]["content"]

            return {"error": False, "content": content, "action": action}

    except Exception as e:
        logger.error(f"Studio action error: {e}")
        return {"error": True, "content": f"Error: {str(e)}"}


# ===================== Notes Service =====================

async def save_note(user_id: str, content: str, source: str = "chat", chat_id: Optional[str] = None) -> NotebookNoteDocument:
    """Save a note."""
    return await NotebookNoteDocument.create(
        user_id=user_id,
        content=content,
        source=source,
        chat_id=chat_id,
    )


async def get_notes(user_id: str) -> List[NotebookNoteDocument]:
    """Get all notes for a user."""
    return await NotebookNoteDocument.find(
        {"user_id": user_id},
        sort=[("created_at", -1)],
        limit=100,
    )


async def delete_note(user_id: str, note_id: str) -> bool:
    """Delete a note."""
    note = await NotebookNoteDocument.find_by_id(note_id)
    if not note or note.user_id != user_id:
        return False
    return await note.delete()
